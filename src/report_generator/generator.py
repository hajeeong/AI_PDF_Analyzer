import os
import json
import logging
from typing import List, Dict, Any, Optional
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np

from ..retrieval.vector_store import LocalVectorStore, create_vector_store
from ..ai_interface.query_processor import QueryProcessor

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class ReportGenerator:
    """
    Generate presentations and summary documents from analyzed annual reports.
    """
    
    def __init__(self, 
                vector_store, 
                query_processor: Optional[QueryProcessor] = None,
                model_name: str = "gpt-3.5-turbo",
                company_name: Optional[str] = None,
                output_dir: Optional[str] = None):
        """
        Initialize the report generator.
        
        Args:
            vector_store: Vector store instance for retrieving report content
            query_processor: Query processor for generating text (optional)
            model_name: LLM model name for query processor (if not provided)
            company_name: Name of the company (optional)
            output_dir: Directory to save generated reports
        """
        self.vector_store = vector_store
        
        # Create query processor if not provided
        if query_processor:
            self.query_processor = query_processor
        else:
            self.query_processor = QueryProcessor(vector_store, model_name=model_name)
        
        # Attempt to extract company name if not provided
        self.company_name = company_name
        if not self.company_name:
            self.company_name = self._extract_company_name()
        
        # Set output directory
        self.output_dir = output_dir or "data/reports"
        os.makedirs(self.output_dir, exist_ok=True)
        
        # Timestamp for generated files
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        logger.info(f"Initialized report generator for {self.company_name or 'unknown company'}")
    
    def _extract_company_name(self) -> Optional[str]:
        """
        Attempt to extract company name from the report content.
        
        Returns:
            Company name or None if not found
        """
        logger.info(f"Attempting to extract company name")
        
        try:
            filename = os.path.basename(self.report_path)
            if "_" in filename:
                name_part = filename.split("_")[0]
                if name_part and len(name_part) > 3:
                    logger.info(f"Using name from filename: {name_part}")
                    return name_part.replace('-', ' ').title()
        except Exception as e:
            logger.error(f"Error extracting name from filename: {str(e)}")
        
        logger.info("Could not extract company name")
        return None
    
    def generate_executive_summary(self, max_length: int = 1000) -> str:
        """
        Generate an executive summary of the annual report.
        
        Args:
            max_length: Maximum length of the summary in words
            
        Returns:
            Executive summary text
        """
        logger.info("Generating executive summary")
        
        # Key sections to include in summary
        key_queries = [
            "What are the key financial highlights and performance of the company in this reporting period?",
            "What is the company's overall business strategy and outlook?",
            "What are the major risks facing the company?",
            "What are the company's main products, services, and markets?",
            "What are the key management changes and corporate governance updates?"
        ]
        
        # Generate responses for each query
        section_responses = {}
        for query in key_queries:
            result = self.query_processor.process_query(query)
            section_responses[query] = result["response"]
        
        # Generate a combined summary
        summary_prompt = f"""
        Create a concise executive summary of {self.company_name or 'the company'}'s annual report.
        The summary should highlight key information from the following sections.
        Use a professional, objective tone appropriate for investors and analysts.
        Keep the summary under {max_length} words.
        
        FINANCIAL HIGHLIGHTS:
        {section_responses[key_queries[0]]}
        
        STRATEGY AND OUTLOOK:
        {section_responses[key_queries[1]]}
        
        KEY RISKS:
        {section_responses[key_queries[2]]}
        
        PRODUCTS AND MARKETS:
        {section_responses[key_queries[3]]}
        
        GOVERNANCE:
        {section_responses[key_queries[4]]}
        
        EXECUTIVE SUMMARY:
        """
        
        try:
            import openai
            response = openai.ChatCompletion.create(
                model=self.query_processor.model_name,
                messages=[{
                    "role": "user",
                    "content": summary_prompt
                }],
                temperature=0.3,
                max_tokens=1500
            )
            
            summary = response.choices[0].message.content.strip()
            logger.info(f"Generated executive summary ({len(summary.split())} words)")
            
            return summary
            
        except Exception as e:
            logger.error(f"Error generating executive summary: {str(e)}")
            # Fallback to a concatenated summary
            fallback_summary = "EXECUTIVE SUMMARY\n\n"
            for i, query in enumerate(key_queries):
                fallback_summary += f"{query.replace('?', '').replace('What are', '').replace('What is', '').strip()}:\n"
                fallback_summary += section_responses[query] + "\n\n"
            
            return fallback_summary
    
    def extract_financial_data(self) -> Dict[str, Any]:
        """
        Extract key financial metrics and data for visualization.
        
        Returns:
            Dictionary containing financial data
        """
        logger.info("Extracting financial data")
        
        # Get financial tables
        table_chunks = self.vector_store.get_table_chunks()
        
        # Extract financial metrics using LLM
        financial_metrics_prompt = """
        Extract key financial metrics from the following tables and text from an annual report.
        For each metric, provide the name, value, previous year value (if available), and percent change.
        Format your response as JSON. If a value is not available, use null.
        
        Example format:
        {
          "revenue": {"value": 1234.5, "previous": 1200.0, "change": 2.88},
          "net_income": {"value": 567.8, "previous": 550.0, "change": 3.24},
          ...
        }
        
        TABLES AND TEXT:
        {text}
        
        JSON FINANCIAL METRICS:
        """
        
        combined_text = "\n\n".join([chunk["text"] for chunk in table_chunks[:5]])
        
        try:
            import openai
            response = openai.ChatCompletion.create(
                model=self.query_processor.model_name,
                messages=[{
                    "role": "user",
                    "content": financial_metrics_prompt.format(text=combined_text)
                }],
                temperature=0.1,
                max_tokens=1000
            )
            
            metrics_text = response.choices[0].message.content.strip()
            
            # Extract JSON from response
            json_match = re.search(r'```json\s*(.*?)\s*```', metrics_text, re.DOTALL)
            if json_match:
                metrics_json = json_match.group(1)
            else:
                metrics_json = metrics_text
            
            # Clean up potential JSON issues
            metrics_json = re.sub(r'```', '', metrics_json)
            metrics_json = metrics_json.strip()
            
            try:
                financial_data = json.loads(metrics_json)
                logger.info(f"Extracted {len(financial_data)} financial metrics")
                return financial_data
            except json.JSONDecodeError as e:
                logger.error(f"Error parsing financial metrics JSON: {str(e)}")
                return {}
            
        except Exception as e:
            logger.error(f"Error extracting financial metrics: {str(e)}")
            return {}
    
    def generate_presentation(self) -> str:
        """
        Generate a PowerPoint presentation summarizing the annual report.
        
        Returns:
            Path to the generated presentation file
        """
        logger.info("Generating presentation")
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        company_name = self.company_name or "Company"
        year = datetime.now().year
        
        title.text = f"{company_name} Annual Report Analysis"
        subtitle.text = f"Executive Briefing\n{year}"
        
        # Add executive summary slide
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        
        # Generate executive summary
        exec_summary = self.generate_executive_summary(max_length=250)
        content.text = exec_summary
        
        # Add financial highlights slide
        financial_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(financial_slide_layout)
        title = slide.shapes.title
        title.text = "Financial Highlights"
        
        # Get financial data
        financial_data = self.extract_financial_data()
        
        # Create a table for financial data
        if financial_data:
            # Create table
            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(3)
            table_rows = min(8, len(financial_data) + 1)  # +1 for header
            table_cols = 4
            
            shape = slide.shapes.add_table(table_rows, table_cols, x, y, cx, cy)
            table = shape.table
            
            # Set header
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Current"
            table.cell(0, 2).text = "Previous"
            table.cell(0, 3).text = "Change (%)"
            
            # Apply header formatting
            for i in range(table_cols):
                cell = table.cell(0, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add data rows
            row_idx = 1
            for metric, data in list(financial_data.items())[:table_rows-1]:
                table.cell(row_idx, 0).text = metric.title()
                
                # Format values
                current = data.get("value")
                previous = data.get("previous")
                change = data.get("change")
                
                table.cell(row_idx, 1).text = f"{current:,.2f}" if current is not None else "N/A"
                table.cell(row_idx, 2).text = f"{previous:,.2f}" if previous is not None else "N/A"
                
                # Format change with color
                if change is not None:
                    change_cell = table.cell(row_idx, 3)
                    change_cell.text = f"{change:+.2f}%"
                    
                    # Color based on positive/negative
                    paragraph = change_cell.text_frame.paragraphs[0]
                    if change > 0:
                        paragraph.font.color.rgb = RGBColor(0, 128, 0)  # Green
                    elif change < 0:
                        paragraph.font.color.rgb = RGBColor(192, 0, 0)  # Red
                else:
                    table.cell(row_idx, 3).text = "N/A"
                
                row_idx += 1
            
            # Create financial chart if possible
            try:
                if len(financial_data) >= 3:
                    # Generate chart
                    chart_path = self._create_financial_chart(financial_data)
                    if chart_path:
                        # Add chart to a new slide
                        chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
                        chart_title = chart_slide.shapes.title
                        chart_title.text = "Financial Performance"
                        
                        # Add chart image
                        chart_slide.shapes.add_picture(
                            chart_path, 
                            Inches(1), 
                            Inches(1.5),
                            width=Inches(8)
                        )
            except Exception as e:
                logger.error(f"Error creating financial chart: {str(e)}")
        
        # Add key strategies slide
        strategy_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(strategy_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Strategy and Outlook"
        
        # Generate strategy content
        strategy_result = self.query_processor.process_query(
            "What is the company's strategy and outlook for the future?")
        content.text = strategy_result["response"]
        
        # Add risk factors slide
        risk_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(risk_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Risk Factors"
        
        # Generate risk content
        risk_query = self.query_processor.process_query(
            "What are the top 5 risk factors for the company?")
        content.text = risk_query["response"]
        
        # Add a slide for ESG/sustainability
        esg_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(esg_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "ESG & Sustainability"
        
        # Generate ESG content
        esg_query = self.query_processor.process_query(
            "What are the company's ESG and sustainability initiatives?")
        content.text = esg_query["response"]
        
        # Save presentation
        company_slug = re.sub(r'[^a-zA-Z0-9]', '_', self.company_name or "company").lower()
        output_path = os.path.join(self.output_dir, 
                                  f"{company_slug}_annual_report_summary_{self.timestamp}.pptx")
        
        prs.save(output_path)
        
        logger.info(f"Saved presentation to {output_path}")
        return output_path
    
    def _create_financial_chart(self, financial_data: Dict[str, Any]) -> Optional[str]:
        """
        Create a financial chart image from extracted data.
        
        Args:
            financial_data: Dictionary of financial metrics
            
        Returns:
            Path to generated chart image or None if failed
        """
        # Filter metrics with both current and previous values
        chart_data = {}
        for metric, data in financial_data.items():
            if data.get("value") is not None and data.get("previous") is not None:
                chart_data[metric] = data
        
        if len(chart_data) < 3:
            logger.warning("Not enough data points for financial chart")
            return None
        
        # Create chart data
        metrics = list(chart_data.keys())
        current_values = [chart_data[m]["value"] for m in metrics]
        previous_values = [chart_data[m]["previous"] for m in metrics]
        
        # Limit to top 5 metrics by current value
        if len(metrics) > 5:
            top_indices = np.argsort(current_values)[-5:]
            metrics = [metrics[i] for i in top_indices]
            current_values = [current_values[i] for i in top_indices]
            previous_values = [previous_values[i] for i in top_indices]
        
        # Create chart
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Set width of bars
        bar_width = 0.35
        indices = np.arange(len(metrics))
        
        # Create bars
        current_bars = ax.bar(indices - bar_width/2, current_values, bar_width, 
                             label='Current Year', color='#4F81BD')
        previous_bars = ax.bar(indices + bar_width/2, previous_values, bar_width,
                              label='Previous Year', color='#C0504D')
        
        # Add labels and title
        ax.set_xlabel('Metrics')
        ax.set_ylabel('Value')
        ax.set_title('Financial Metrics Comparison')
        ax.set_xticks(indices)
        
        # Format metric names for display
        display_metrics = [m.replace('_', ' ').title() for m in metrics]
        display_metrics = [m[:15] + '...' if len(m) > 15 else m for m in display_metrics]
        ax.set_xticklabels(display_metrics, rotation=45, ha='right')
        
        # Add value labels on bars
        for bar in current_bars:
            height = bar.get_height()
            ax.annotate(f'{height:,.1f}',
                       xy=(bar.get_x() + bar.get_width() / 2, height),
                       xytext=(0, 3),  # 3 points vertical offset
                       textcoords="offset points",
                       ha='center', va='bottom', fontsize=8)
        
        for bar in previous_bars:
            height = bar.get_height()
            ax.annotate(f'{height:,.1f}',
                       xy=(bar.get_x() + bar.get_width() / 2, height),
                       xytext=(0, 3),  # 3 points vertical offset
                       textcoords="offset points",
                       ha='center', va='bottom', fontsize=8)
        
        # Add legend
        ax.legend()
        
        # Save figure
        plt.tight_layout()
        chart_path = os.path.join(self.output_dir, f"financial_chart_{self.timestamp}.png")
        plt.savefig(chart_path, dpi=300)
        plt.close()
        
        logger.info(f"Created financial chart at {chart_path}")
        return chart_path
    
    def generate_text_report(self) -> str:
        """
        Generate a text-based report summarizing the annual report.
        
        Returns:
            Path to the generated report file
        """
        logger.info("Generating text report")
        
        # Sections to include in the report
        sections = [
            {
                "title": "EXECUTIVE SUMMARY",
                "query": None,  # Will use generate_executive_summary
                "special": "executive_summary"
            },
            {
                "title": "FINANCIAL HIGHLIGHTS",
                "query": "What are the key financial highlights and performance metrics? Include specific numbers and comparisons to previous periods.",
                "special": None
            },
            {
                "title": "BUSINESS OVERVIEW",
                "query": "Provide an overview of the company's business, including its main products, services, and markets.",
                "special": None
            },
            {
                "title": "STRATEGY AND OUTLOOK",
                "query": "What is the company's strategy and outlook for the future?",
                "special": None
            },
            {
                "title": "RISK FACTORS",
                "query": "What are the top risk factors facing the company? Summarize each key risk.",
                "special": None
            },
            {
                "title": "OPERATIONS REVIEW",
                "query": "How did the company's operations perform during this period? Discuss key operational highlights and challenges.",
                "special": None
            },
            {
                "title": "GOVERNANCE AND LEADERSHIP",
                "query": "Provide information about the company's governance structure, board of directors, and any leadership changes.",
                "special": None
            },
            {
                "title": "ESG AND SUSTAINABILITY",
                "query": "What are the company's ESG (Environmental, Social, Governance) initiatives and sustainability efforts?",
                "special": None
            },
            {
                "title": "CONCLUSION AND RECOMMENDATIONS",
                "query": None,
                "special": "conclusion"
            }
        ]
        
        # Generate content for each section
        report_content = f"# {self.company_name or 'COMPANY'} ANNUAL REPORT ANALYSIS\n\n"
        report_content += f"*Generated on {datetime.now().strftime('%B %d, %Y')}*\n\n"
        
        for section in sections:
            report_content += f"## {section['title']}\n\n"
            
            if section["special"] == "executive_summary":
                # Use executive summary generator
                content = self.generate_executive_summary()
            elif section["special"] == "conclusion":
                # Generate conclusion based on previous sections
                conclusion_query = "Based on all the information in this annual report, what are the key takeaways and potential recommendations for investors?"
                result = self.query_processor.process_query(conclusion_query)
                content = result["response"]
            else:
                # Use standard query
                result = self.query_processor.process_query(section["query"])
                content = result["response"]
            
            report_content += f"{content}\n\n"
        
        # Add financial data if available
        financial_data = self.extract_financial_data()
        if financial_data:
            report_content += "## FINANCIAL DATA TABLE\n\n"
            report_content += "| Metric | Current | Previous | Change (%) |\n"
            report_content += "|--------|---------|----------|------------|\n"
            
            for metric, data in financial_data.items():
                current = data.get("value", "N/A")
                previous = data.get("previous", "N/A")
                change = data.get("change", "N/A")
                
                if current != "N/A":
                    current = f"{current:,.2f}"
                if previous != "N/A":
                    previous = f"{previous:,.2f}"
                if change != "N/A":
                    change = f"{change:+.2f}%"
                
                report_content += f"| {metric.title()} | {current} | {previous} | {change} |\n"
            
            report_content += "\n"
        
        # Save report
        company_slug = re.sub(r'[^a-zA-Z0-9]', '_', self.company_name or "company").lower()
        output_path = os.path.join(self.output_dir, 
                                  f"{company_slug}_annual_report_summary_{self.timestamp}.md")
        
        with open(output_path, 'w') as f:
            f.write(report_content)
        
        logger.info(f"Saved text report to {output_path}")
        return output_path


# Example usage
if __name__ == "__main__":
    import sys
    from vector_store import LocalVectorStore
    
    if len(sys.argv) < 2:
        print("Usage: python generator.py <vector_store_dir> [company_name]")
        sys.exit(1)
    
    vector_store_dir = sys.argv[1]
    company_name = sys.argv[2] if len(sys.argv) > 2 else None
    
    vector_store = LocalVectorStore(vector_store_dir)
    generator = ReportGenerator(vector_store, company_name=company_name)
    
    # Generate reports
    text_report_path = generator.generate_text_report()
    ppt_path = generator.generate_presentation()
    
    print(f"Generated text report: {text_report_path}")
    print(f"Generated presentation: {ppt_path}")